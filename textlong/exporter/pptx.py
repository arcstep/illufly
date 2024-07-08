from ..utils import raise_not_install
from ..importer import load_markdown
from pptx import Presentation
from pptx.util import Inches, Pt

def export_pptx(input_file: str, output_file: str = None):
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError as e:
        raise_not_install("please install python-pptx: `pip install python-pptx`")

    if not input_file:
        return False

    md = load_markdown(input_file)

    if not md.documents:
        return False

    if not output_file:
        output_file = input_file + ".pptx"

    prs = Presentation()
    mode, segments = md.get_todo_documents("segment")
    for docs in segments:
        slide_layout = prs.slide_layouts[1]  # 0是标题幻灯片，1是标题和内容
        slide = prs.slides.add_slide(slide_layout)
        title_set = False

        for doc in docs:
            print('type >>', doc.metadata['type'])
            if doc.metadata['type'] in ['front-matter', 'blank_line']:
                continue
            elif doc.metadata['type'] == 'heading' and not title_set:
                title = slide.shapes.title
                title.text = doc.page_content
                title_set = True
                
                print("heading >>", doc.page_content)
            elif doc.page_content:
                if not title_set:  # 如果还没有设置标题，使用第一个内容作为标题
                    title = slide.shapes.title
                    title.text = doc.page_content
                    title_set = True

                    print("not title set >>", doc.page_content)
                else:
                    content_placeholder = slide.placeholders[1]
                    content_text_frame = content_placeholder.text_frame
                    if content_text_frame.text:  # 如果已经有文本，添加新行
                        content_text_frame.add_paragraph().text = doc.page_content
                    else:  # 如果没有文本，直接设置文本
                        content_text_frame.text = doc.page_content
                
                    print("text >>", doc.page_content)
                    print(doc.page_content)

    # 保存PPT文件
    prs.save(output_file)

    return True